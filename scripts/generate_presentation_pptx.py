"""Generate the 6-slide PowerPoint for the AI in FinTech project.

Slide plan (mirrors mhmd/FinTech_Fast_Speaking_5_Minute_Script_Larger.pdf):
  1. Title + team (the six names)
  2. Business problem & project aim
  3. Dataset & inputs
  4. Methodology (binning, WOE, logistic regression, validation)
  5. Portfolio decision (accept / reject)
  6. Expected loss, risk pricing & conclusion

Run from the repo root:
    python3 scripts/generate_presentation_pptx.py
Output: mhmd/fintech_presentation_6slides.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
GOLD = RGBColor(0xE6, 0xB3, 0x3E)
LIGHT = RGBColor(0xF4, 0xF6, 0xFB)
TEXT = RGBColor(0x1A, 0x1F, 0x2E)
MUTED = RGBColor(0x55, 0x5F, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1F, 0xA8, 0x6B)
RED = RGBColor(0xD0, 0x3A, 0x3A)


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=18,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    font="Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items,
    *,
    size=16,
    color=TEXT,
    bullet_color=ACCENT,
    line_spacing=1.25,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        bullet_run = p.add_run()
        bullet_run.text = "\u25CF  "
        bullet_run.font.name = "Calibri"
        bullet_run.font.size = Pt(size)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = bullet_color
        text_run = p.add_run()
        text_run.text = item
        text_run.font.name = "Calibri"
        text_run.font.size = Pt(size)
        text_run.font.color.rgb = color
    return box


def add_footer(slide, page_number, total_pages, slide_w, slide_h):
    add_rect(slide, Emu(0), slide_h - Inches(0.35), slide_w, Inches(0.35), NAVY)
    add_text(
        slide,
        Inches(0.5),
        slide_h - Inches(0.33),
        Inches(8),
        Inches(0.3),
        "AI in FinTech  |  Portfolio Optimisation & Interest Rate Pricing",
        size=10,
        color=WHITE,
    )
    add_text(
        slide,
        slide_w - Inches(2.0),
        slide_h - Inches(0.33),
        Inches(1.5),
        Inches(0.3),
        f"{page_number} / {total_pages}",
        size=10,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
    )


def add_slide_header(slide, slide_w, kicker, title):
    add_rect(slide, Emu(0), Emu(0), slide_w, Inches(0.18), ACCENT)
    add_text(
        slide,
        Inches(0.5),
        Inches(0.32),
        Inches(12),
        Inches(0.35),
        kicker.upper(),
        size=12,
        bold=True,
        color=ACCENT,
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(0.6),
        Inches(12),
        Inches(0.7),
        title,
        size=30,
        bold=True,
        color=NAVY,
    )
    add_rect(
        slide, Inches(0.5), Inches(1.32), Inches(0.6), Inches(0.05), GOLD
    )


def build_title_slide(prs, slide_w, slide_h):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Emu(0), Emu(0), slide_w, slide_h, NAVY)

    add_rect(
        slide, Inches(0.5), Inches(0.5), Inches(2.4), Inches(0.45), ACCENT
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(0.53),
        Inches(2.4),
        Inches(0.4),
        "AI IN FINTECH  |  PROJECT 1",
        size=12,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    add_text(
        slide,
        Inches(0.5),
        Inches(1.4),
        Inches(12.3),
        Inches(1.2),
        "Portfolio Optimisation &",
        size=44,
        bold=True,
        color=WHITE,
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(2.05),
        Inches(12.3),
        Inches(1.2),
        "Interest Rate Pricing",
        size=44,
        bold=True,
        color=WHITE,
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(2.85),
        Inches(12.3),
        Inches(0.6),
        "A Credit Scoring Model Approach",
        size=22,
        color=GOLD,
    )

    add_rect(slide, Inches(0.5), Inches(3.55), Inches(0.6), Inches(0.05), GOLD)

    add_text(
        slide,
        Inches(0.5),
        Inches(3.75),
        Inches(6),
        Inches(0.4),
        "TEAM MEMBERS",
        size=13,
        bold=True,
        color=GOLD,
    )

    members = [
        "Mhamad Naser Dine",
        "Ali Srour",
        "Marwa Zeineddine",
        "Maroun Mashaalany",
        "Lyn Khalil",
        "Ghandour Ghandour",
    ]

    col_w = Inches(3.95)
    row_h = Inches(0.55)
    gap_x = Inches(0.15)
    start_x = Inches(0.5)
    start_y = Inches(4.2)
    for i, name in enumerate(members):
        col = i % 3
        row = i // 3
        x = start_x + (col_w + gap_x) * col
        y = start_y + (row_h + Inches(0.15)) * row
        card = add_rect(slide, x, y, col_w, row_h, RGBColor(0x14, 0x2A, 0x4A))
        card.line.color.rgb = ACCENT
        card.line.width = Emu(12700)
        add_rect(slide, x, y, Inches(0.08), row_h, GOLD)
        add_text(
            slide,
            x + Inches(0.25),
            y + Inches(0.12),
            col_w - Inches(0.3),
            Inches(0.4),
            name,
            size=15,
            bold=True,
            color=WHITE,
        )

    add_rect(
        slide,
        Inches(0.5),
        Inches(6.4),
        Inches(12.3),
        Inches(0.04),
        ACCENT,
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(6.55),
        Inches(7),
        Inches(0.35),
        "Course: AI in FinTech   |   Academic Year 2025 / 2026",
        size=13,
        color=WHITE,
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(6.85),
        Inches(7),
        Inches(0.35),
        "Software: MATLAB Credit Scorecard Workflow",
        size=12,
        color=GOLD,
    )


def build_problem_slide(prs, slide_w, slide_h):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Emu(0), Emu(0), slide_w, slide_h, LIGHT)
    add_slide_header(slide, slide_w, "02  |  Context", "Business Problem & Project Aim")

    add_text(
        slide,
        Inches(0.5),
        Inches(1.55),
        Inches(12.3),
        Inches(0.5),
        "The bank has to decide who gets a loan, and at what rate. We replace gut feel with a data-driven scorecard.",
        size=15,
        color=MUTED,
    )

    card_y = Inches(2.25)
    card_h = Inches(2.4)
    card_w = Inches(4.05)
    gap = Inches(0.15)

    cards = [
        ("THE PROBLEM", "If the bank accepts everyone, defaults destroy profit. If it rejects too many, it loses business.", RED),
        ("THE AIM", "Build a scoring model that ranks clients, decides accept/reject, and prices the minimum interest rate.", ACCENT),
        ("THE OUTPUT", "Score 0\u2013100, probability of default, decision, expected loss and break-even rate per client.", GREEN),
    ]
    for i, (title, body, accent) in enumerate(cards):
        x = Inches(0.5) + (card_w + gap) * i
        card = add_rect(slide, x, card_y, card_w, card_h, WHITE)
        card.line.color.rgb = RGBColor(0xE2, 0xE6, 0xF0)
        add_rect(slide, x, card_y, card_w, Inches(0.18), accent)
        add_text(
            slide,
            x + Inches(0.25),
            card_y + Inches(0.32),
            card_w - Inches(0.5),
            Inches(0.4),
            title,
            size=13,
            bold=True,
            color=accent,
        )
        add_text(
            slide,
            x + Inches(0.25),
            card_y + Inches(0.85),
            card_w - Inches(0.5),
            card_h - Inches(1),
            body,
            size=14,
            color=TEXT,
        )

    add_rect(slide, Inches(0.5), Inches(4.95), Inches(12.3), Inches(1.6), NAVY)
    add_text(
        slide,
        Inches(0.75),
        Inches(5.1),
        Inches(11),
        Inches(0.4),
        "DECISION RULE",
        size=12,
        bold=True,
        color=GOLD,
    )
    add_text(
        slide,
        Inches(0.75),
        Inches(5.4),
        Inches(11),
        Inches(0.6),
        "Accept if Score \u2265 48.2032   \u2192   else Reject",
        size=24,
        bold=True,
        color=WHITE,
    )
    add_text(
        slide,
        Inches(0.75),
        Inches(5.95),
        Inches(11.5),
        Inches(0.5),
        "Cutoff chosen by the Kolmogorov\u2013Smirnov statistic on the historical sample.",
        size=13,
        color=RGBColor(0xC9, 0xD3, 0xE6),
    )


def build_dataset_slide(prs, slide_w, slide_h):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Emu(0), Emu(0), slide_w, slide_h, LIGHT)
    add_slide_header(slide, slide_w, "03  |  Data", "Dataset & Inputs")

    add_text(
        slide,
        Inches(0.5),
        Inches(1.55),
        Inches(12.3),
        Inches(0.5),
        "Source file: DataProjScoreCard.xlsx \u2014 two sheets, four predictors, one binary target.",
        size=15,
        color=MUTED,
    )

    box_w = Inches(6.05)
    box_h = Inches(2.4)
    box_y = Inches(2.25)

    hist = add_rect(slide, Inches(0.5), box_y, box_w, box_h, WHITE)
    hist.line.color.rgb = RGBColor(0xE2, 0xE6, 0xF0)
    add_rect(slide, Inches(0.5), box_y, box_w, Inches(0.45), ACCENT)
    add_text(
        slide,
        Inches(0.7),
        box_y + Inches(0.08),
        box_w,
        Inches(0.4),
        "HistoricalData  \u2014  trains & validates the model",
        size=14,
        bold=True,
        color=WHITE,
    )
    add_text(
        slide,
        Inches(0.7),
        box_y + Inches(0.65),
        Inches(2.5),
        Inches(0.6),
        "500",
        size=44,
        bold=True,
        color=NAVY,
    )
    add_text(
        slide,
        Inches(0.7),
        box_y + Inches(1.45),
        Inches(5.5),
        Inches(0.4),
        "past clients with known Default outcome",
        size=13,
        color=TEXT,
    )
    add_text(
        slide,
        Inches(0.7),
        box_y + Inches(1.8),
        Inches(5.5),
        Inches(0.4),
        "320 good (64%)   \u2022   180 bad (36%)",
        size=13,
        color=MUTED,
    )

    port_x = Inches(0.5) + box_w + Inches(0.2)
    port = add_rect(slide, port_x, box_y, box_w, box_h, WHITE)
    port.line.color.rgb = RGBColor(0xE2, 0xE6, 0xF0)
    add_rect(slide, port_x, box_y, box_w, Inches(0.45), GOLD)
    add_text(
        slide,
        port_x + Inches(0.2),
        box_y + Inches(0.08),
        box_w,
        Inches(0.4),
        "ActualPortfolioData  \u2014  receives the decision",
        size=14,
        bold=True,
        color=NAVY,
    )
    add_text(
        slide,
        port_x + Inches(0.2),
        box_y + Inches(0.65),
        Inches(2.5),
        Inches(0.6),
        "20",
        size=44,
        bold=True,
        color=NAVY,
    )
    add_text(
        slide,
        port_x + Inches(0.2),
        box_y + Inches(1.45),
        Inches(5.5),
        Inches(0.4),
        "new applicants \u2014 outcome unknown",
        size=13,
        color=TEXT,
    )
    add_text(
        slide,
        port_x + Inches(0.2),
        box_y + Inches(1.8),
        Inches(5.5),
        Inches(0.4),
        "Each requests a loan of $100,000",
        size=13,
        color=MUTED,
    )

    add_text(
        slide,
        Inches(0.5),
        Inches(4.85),
        Inches(12),
        Inches(0.4),
        "FOUR PREDICTORS",
        size=12,
        bold=True,
        color=ACCENT,
    )

    preds = [
        ("Age", "numeric", "client age in years"),
        ("Income", "numeric", "annual income in $"),
        ("Residential Status", "categorical", "renter / homeowner"),
        ("Employment Status", "categorical", "employed / other"),
    ]
    pw = Inches(2.95)
    py = Inches(5.25)
    for i, (name, kind, desc) in enumerate(preds):
        x = Inches(0.5) + (pw + Inches(0.1)) * i
        pcard = add_rect(slide, x, py, pw, Inches(1.3), WHITE)
        pcard.line.color.rgb = RGBColor(0xE2, 0xE6, 0xF0)
        add_rect(slide, x, py, Inches(0.08), Inches(1.3), GOLD)
        add_text(
            slide,
            x + Inches(0.2),
            py + Inches(0.15),
            pw - Inches(0.3),
            Inches(0.35),
            name,
            size=15,
            bold=True,
            color=NAVY,
        )
        add_text(
            slide,
            x + Inches(0.2),
            py + Inches(0.5),
            pw - Inches(0.3),
            Inches(0.3),
            kind.upper(),
            size=10,
            bold=True,
            color=ACCENT,
        )
        add_text(
            slide,
            x + Inches(0.2),
            py + Inches(0.8),
            pw - Inches(0.3),
            Inches(0.4),
            desc,
            size=12,
            color=MUTED,
        )


def build_method_slide(prs, slide_w, slide_h):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Emu(0), Emu(0), slide_w, slide_h, LIGHT)
    add_slide_header(slide, slide_w, "04  |  Method", "Scorecard Workflow & Validation")

    add_text(
        slide,
        Inches(0.5),
        Inches(1.55),
        Inches(12.3),
        Inches(0.5),
        "Five-step MATLAB pipeline turns raw client data into an interpretable 0\u2013100 score.",
        size=15,
        color=MUTED,
    )

    steps = ["Binning", "Weight of Evidence", "Logistic Regression", "Score 0\u2013100", "ROC Validation"]
    sx = Inches(0.5)
    sy = Inches(2.3)
    sw = Inches(2.34)
    sh = Inches(0.8)
    for i, s in enumerate(steps):
        x = sx + (sw + Inches(0.05)) * i
        chip = add_rect(slide, x, sy, sw, sh, NAVY)
        add_text(
            slide,
            x,
            sy + Inches(0.15),
            sw,
            Inches(0.3),
            f"STEP {i + 1}",
            size=10,
            bold=True,
            color=GOLD,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            x,
            sy + Inches(0.42),
            sw,
            Inches(0.35),
            s,
            size=13,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                x + sw - Inches(0.02),
                sy + Inches(0.28),
                Inches(0.18),
                Inches(0.24),
            )
            arrow.rotation = 90
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()

    iv_x = Inches(0.5)
    iv_y = Inches(3.4)
    iv_w = Inches(6.4)
    iv_h = Inches(3.15)
    iv_card = add_rect(slide, iv_x, iv_y, iv_w, iv_h, WHITE)
    iv_card.line.color.rgb = RGBColor(0xE2, 0xE6, 0xF0)
    add_text(
        slide,
        iv_x + Inches(0.25),
        iv_y + Inches(0.2),
        iv_w - Inches(0.5),
        Inches(0.4),
        "PREDICTOR STRENGTH (Information Value)",
        size=12,
        bold=True,
        color=ACCENT,
    )

    ivs = [
        ("Income", 0.2355, "Strong"),
        ("Age", 0.1879, "Medium"),
        ("Employment Status", 0.1143, "Medium"),
        ("Residential Status", 0.0417, "Weak"),
    ]
    max_iv = max(v for _, v, _ in ivs)
    bar_left = iv_x + Inches(2.1)
    bar_max_w = iv_w - Inches(3.3)
    row_y = iv_y + Inches(0.75)
    row_step = Inches(0.55)
    for i, (name, value, strength) in enumerate(ivs):
        y = row_y + row_step * i
        add_text(slide, iv_x + Inches(0.25), y, Inches(1.85), Inches(0.4), name, size=12, bold=True, color=TEXT)
        bg = add_rect(slide, bar_left, y + Inches(0.08), bar_max_w, Inches(0.22), RGBColor(0xEC, 0xEF, 0xF5))
        w = Emu(int(bar_max_w * (value / max_iv)))
        add_rect(slide, bar_left, y + Inches(0.08), w, Inches(0.22), ACCENT)
        add_text(
            slide,
            bar_left + bar_max_w + Inches(0.05),
            y,
            Inches(0.9),
            Inches(0.4),
            f"{value:.4f}",
            size=11,
            bold=True,
            color=NAVY,
        )

    val_x = Inches(7.05)
    val_y = Inches(3.4)
    val_w = Inches(5.75)
    val_card = add_rect(slide, val_x, val_y, val_w, iv_h, NAVY)
    add_text(
        slide,
        val_x + Inches(0.25),
        val_y + Inches(0.2),
        val_w - Inches(0.5),
        Inches(0.4),
        "MODEL VALIDATION",
        size=12,
        bold=True,
        color=GOLD,
    )

    metrics = [
        ("AUC", "0.6775", "Area under ROC curve"),
        ("KS", "0.2892", "Kolmogorov\u2013Smirnov"),
        ("Cutoff", "48.20", "Optimal score threshold"),
    ]
    mh = Inches(0.78)
    for i, (label, value, sub) in enumerate(metrics):
        y = val_y + Inches(0.75) + (mh + Inches(0.04)) * i
        add_rect(slide, val_x + Inches(0.25), y, val_w - Inches(0.5), mh, RGBColor(0x14, 0x2A, 0x4A))
        add_text(
            slide,
            val_x + Inches(0.45),
            y + Inches(0.15),
            Inches(1.5),
            Inches(0.5),
            label,
            size=12,
            bold=True,
            color=GOLD,
        )
        add_text(
            slide,
            val_x + Inches(1.8),
            y + Inches(0.1),
            Inches(1.8),
            Inches(0.55),
            value,
            size=22,
            bold=True,
            color=WHITE,
        )
        add_text(
            slide,
            val_x + Inches(3.4),
            y + Inches(0.22),
            Inches(2.1),
            Inches(0.4),
            sub,
            size=10,
            color=RGBColor(0xC9, 0xD3, 0xE6),
        )


def build_decision_slide(prs, slide_w, slide_h):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Emu(0), Emu(0), slide_w, slide_h, LIGHT)
    add_slide_header(slide, slide_w, "05  |  Decision", "Portfolio Accept / Reject")

    add_text(
        slide,
        Inches(0.5),
        Inches(1.55),
        Inches(12.3),
        Inches(0.5),
        "Rule applied: Accept if Score \u2265 48.2032 \u2014 systematic, evidence-based, defensible.",
        size=15,
        color=MUTED,
    )

    kpis = [
        ("14", "ACCEPTED", "out of 20 applicants", GREEN),
        ("6", "REJECTED", "highest-risk profiles", RED),
        ("$1.4M", "EXPOSURE", "total accepted EAD", ACCENT),
        ("23.09%", "AVG PD", "of accepted portfolio", GOLD),
    ]
    kw = Inches(2.95)
    kh = Inches(1.6)
    ky = Inches(2.25)
    for i, (big, label, sub, color) in enumerate(kpis):
        x = Inches(0.5) + (kw + Inches(0.1)) * i
        card = add_rect(slide, x, ky, kw, kh, WHITE)
        card.line.color.rgb = RGBColor(0xE2, 0xE6, 0xF0)
        add_rect(slide, x, ky, kw, Inches(0.12), color)
        add_text(
            slide,
            x + Inches(0.25),
            ky + Inches(0.3),
            kw - Inches(0.5),
            Inches(0.7),
            big,
            size=34,
            bold=True,
            color=NAVY,
        )
        add_text(
            slide,
            x + Inches(0.25),
            ky + Inches(0.95),
            kw - Inches(0.5),
            Inches(0.3),
            label,
            size=11,
            bold=True,
            color=color,
        )
        add_text(
            slide,
            x + Inches(0.25),
            ky + Inches(1.22),
            kw - Inches(0.5),
            Inches(0.3),
            sub,
            size=11,
            color=MUTED,
        )

    band_y = Inches(4.1)
    band_h = Inches(2.45)
    band_card = add_rect(slide, Inches(0.5), band_y, Inches(12.3), band_h, WHITE)
    band_card.line.color.rgb = RGBColor(0xE2, 0xE6, 0xF0)
    add_text(
        slide,
        Inches(0.75),
        band_y + Inches(0.2),
        Inches(12),
        Inches(0.4),
        "RISK BAND DISTRIBUTION",
        size=12,
        bold=True,
        color=ACCENT,
    )

    bands = [
        ("Low Risk", 5, "Score > 80", GREEN, "Clients 4, 5, 6, 7, 18"),
        ("Medium Risk", 9, "50 \u2264 Score \u2264 80", GOLD, "Mid-band accepted clients"),
        ("High Risk", 6, "Score < 50", RED, "All rejected applicants"),
    ]
    by = band_y + Inches(0.7)
    bw = (Inches(12.3) - Inches(0.5) * 2 - Inches(0.3)) / 3
    for i, (name, count, rule, color, members) in enumerate(bands):
        x = Inches(0.75) + (bw + Inches(0.15)) * i
        add_rect(slide, x, by, bw, Inches(1.55), LIGHT)
        add_rect(slide, x, by, Inches(0.08), Inches(1.55), color)
        add_text(slide, x + Inches(0.2), by + Inches(0.12), bw - Inches(0.3), Inches(0.35), name, size=14, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.2), by + Inches(0.45), Inches(0.8), Inches(0.55), str(count), size=28, bold=True, color=color)
        add_text(slide, x + Inches(1.1), by + Inches(0.6), bw - Inches(1.2), Inches(0.3), "clients", size=11, color=MUTED)
        add_text(slide, x + Inches(0.2), by + Inches(1.0), bw - Inches(0.3), Inches(0.3), rule, size=11, bold=True, color=color)
        add_text(slide, x + Inches(0.2), by + Inches(1.25), bw - Inches(0.3), Inches(0.3), members, size=10, color=MUTED)


def build_conclusion_slide(prs, slide_w, slide_h):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Emu(0), Emu(0), slide_w, slide_h, LIGHT)
    add_slide_header(slide, slide_w, "06  |  Conclusion", "Expected Loss, Pricing & Wrap-up")

    formula_y = Inches(1.55)
    formula = add_rect(slide, Inches(0.5), formula_y, Inches(12.3), Inches(0.9), NAVY)
    add_text(
        slide,
        Inches(0.75),
        formula_y + Inches(0.1),
        Inches(12),
        Inches(0.3),
        "EXPECTED LOSS & MINIMUM BREAK-EVEN RATE",
        size=11,
        bold=True,
        color=GOLD,
    )
    add_text(
        slide,
        Inches(0.75),
        formula_y + Inches(0.4),
        Inches(12),
        Inches(0.5),
        "EL = PD \u00D7 LGD \u00D7 EAD     |     r\u2098\u1D62\u2099 = PD \u00D7 LGD     (LGD = 40%, EAD = $100,000)",
        size=16,
        bold=True,
        color=WHITE,
    )

    metrics = [
        ("Total Expected Loss", "$129,329.63", "across 14 accepted loans", RED),
        ("Avg Break-even Rate", "9.24%", "credit-risk floor", ACCENT),
        ("Safest Client Rate", "5.11%", "clients 4 & 5 (PD 12.79%)", GREEN),
        ("Riskiest Accepted", "13.40%", "client 8 (PD 33.50%)", GOLD),
    ]
    mx = Inches(0.5)
    my = Inches(2.75)
    mw = Inches(2.95)
    mh = Inches(1.7)
    for i, (label, value, sub, color) in enumerate(metrics):
        x = mx + (mw + Inches(0.1)) * i
        card = add_rect(slide, x, my, mw, mh, WHITE)
        card.line.color.rgb = RGBColor(0xE2, 0xE6, 0xF0)
        add_rect(slide, x, my, Inches(0.08), mh, color)
        add_text(slide, x + Inches(0.25), my + Inches(0.2), mw - Inches(0.4), Inches(0.4), label.upper(), size=10, bold=True, color=color)
        add_text(slide, x + Inches(0.25), my + Inches(0.55), mw - Inches(0.4), Inches(0.7), value, size=26, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.25), my + Inches(1.25), mw - Inches(0.4), Inches(0.4), sub, size=11, color=MUTED)

    take_y = Inches(4.65)
    take_card = add_rect(slide, Inches(0.5), take_y, Inches(7.8), Inches(1.9), WHITE)
    take_card.line.color.rgb = RGBColor(0xE2, 0xE6, 0xF0)
    add_text(
        slide,
        Inches(0.75),
        take_y + Inches(0.2),
        Inches(7.3),
        Inches(0.4),
        "KEY TAKEAWAYS",
        size=12,
        bold=True,
        color=ACCENT,
    )
    add_bullets(
        slide,
        Inches(0.75),
        take_y + Inches(0.6),
        Inches(7.3),
        Inches(1.3),
        [
            "Score \u2192 PD \u2192 Decision \u2192 Expected Loss \u2192 Risk-based Price",
            "Transparent, interpretable, regulator-friendly scorecard",
            "Rate above 9.24% generates positive expected margin for the bank",
        ],
        size=13,
    )

    thanks = add_rect(slide, Inches(8.45), take_y, Inches(4.35), Inches(1.9), NAVY)
    add_text(
        slide,
        Inches(8.7),
        take_y + Inches(0.35),
        Inches(4),
        Inches(0.5),
        "Thank you",
        size=28,
        bold=True,
        color=WHITE,
    )
    add_text(
        slide,
        Inches(8.7),
        take_y + Inches(0.95),
        Inches(4),
        Inches(0.4),
        "Questions & discussion",
        size=14,
        color=GOLD,
    )
    add_rect(slide, Inches(8.7), take_y + Inches(1.35), Inches(0.5), Inches(0.05), GOLD)
    add_text(
        slide,
        Inches(8.7),
        take_y + Inches(1.45),
        Inches(4),
        Inches(0.35),
        "AI in FinTech \u2014 Project 1",
        size=11,
        color=RGBColor(0xC9, 0xD3, 0xE6),
    )


def main() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    builders = [
        build_title_slide,
        build_problem_slide,
        build_dataset_slide,
        build_method_slide,
        build_decision_slide,
        build_conclusion_slide,
    ]

    total = len(builders)
    for i, builder in enumerate(builders, start=1):
        builder(prs, slide_w, slide_h)
        if i > 1:
            add_footer(prs.slides[i - 1], i, total, slide_w, slide_h)

    out_dir = Path(__file__).resolve().parent.parent / "mhmd"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "fintech_presentation_6slides.pptx"
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = main()
    print(f"Saved: {path}")
